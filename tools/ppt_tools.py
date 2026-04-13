from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from sandbox import resolve_path
import os
import json
import textwrap

# ─────────────────────────────────────────────
# Constantes de diseño
# ─────────────────────────────────────────────
COLOR_BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # Azul noche oscuro
COLOR_BG_ACCENT = RGBColor(0x16, 0x21, 0x3E)   # Azul medianoche
COLOR_BG_REF    = RGBColor(0x0F, 0x0F, 0x1A)   # Casi negro para referencias
COLOR_TITLE     = RGBColor(0xE0, 0xE0, 0xFF)   # Blanco azulado
COLOR_BODY      = RGBColor(0xD0, 0xD0, 0xD0)   # Gris claro
COLOR_ACCENT    = RGBColor(0x4F, 0xA3, 0xFF)   # Azul brillante para acentos
COLOR_REF_TEXT  = RGBColor(0xAA, 0xAA, 0xCC)   # Gris violáceo para refs

FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"

MAX_VIÑETAS_POR_SLIDE = 5     # Si hay más, se divide en slides extra
MAX_CHARS_VIÑETA      = 120   # Truncar viñetas largas
MAX_REFS_POR_SLIDE    = 6     # Referencias por slide de referencias
MAX_FILAS_TABLA_SLIDE = 8     # Filas de datos por slide de tabla (sin encabezado)


# ─────────────────────────────────────────────
# Helpers visuales
# ─────────────────────────────────────────────

def _set_slide_background(slide, color: RGBColor):
    """Rellena el fondo del slide con un color sólido."""
    from pptx.oxml.ns import qn
    from lxml import etree

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_textbox(slide, titulo: str, prs: Presentation):
    """Añade el título del slide con estilo definido."""
    w, h = prs.slide_width, prs.slide_height
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        w - Inches(1), Inches(1.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = titulo
    run.font.name = FONT_TITLE
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLOR_TITLE


def _add_body_textbox(slide, lines: list, prs: Presentation, top_offset: float = 1.6):
    """Añade las viñetas/cuerpo del slide."""
    w, h = prs.slide_width, prs.slide_height
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(top_offset),
        w - Inches(1.1), h - Inches(top_offset + 0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        # Truncar viñetas muy largas
        line = line if len(line) <= MAX_CHARS_VIÑETA else line[:MAX_CHARS_VIÑETA].rstrip() + "…"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"▸  {line}"
        run.font.name = FONT_BODY
        run.font.size = Pt(17)
        run.font.color.rgb = COLOR_BODY


def _add_accent_bar(slide, prs: Presentation):
    """Añade una barra de color en la parte izquierda del slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(0.12), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()


def _blank_slide(prs: Presentation):
    """Crea un slide en blanco (sin placeholders)."""
    blank_layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(blank_layout)


# ─────────────────────────────────────────────
# Builders por tipo
# ─────────────────────────────────────────────

def _build_portada(prs: Presentation, titulo: str, subtitulo: str):
    slide = _blank_slide(prs)
    _set_slide_background(slide, COLOR_BG_DARK)
    _add_accent_bar(slide, prs)

    w, h = prs.slide_width, prs.slide_height

    # Título grande centrado verticalmente
    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), w - Inches(1.4), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = titulo
    run.font.name = FONT_TITLE
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = COLOR_TITLE

    # Subtítulo
    txSub = slide.shapes.add_textbox(Inches(0.9), Inches(4.4), w - Inches(1.4), Inches(1.0))
    tf2 = txSub.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = subtitulo
    run2.font.name = FONT_BODY
    run2.font.size = Pt(18)
    run2.font.color.rgb = COLOR_ACCENT


def _build_cita(prs: Presentation, texto: str):
    slide = _blank_slide(prs)
    _set_slide_background(slide, COLOR_BG_ACCENT)

    w, h = prs.slide_width, prs.slide_height
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), w - Inches(2.0), h - Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f'❝  {texto}  ❞'
    run.font.name = FONT_BODY
    run.font.size = Pt(22)
    run.font.italic = True
    run.font.color.rgb = COLOR_TITLE


def _build_contenido(prs: Presentation, titulo: str, viñetas: list):
    """
    Si hay más de MAX_VIÑETAS_POR_SLIDE viñetas, crea múltiples slides automáticamente.
    """
    chunks = [viñetas[i:i+MAX_VIÑETAS_POR_SLIDE]
              for i in range(0, max(len(viñetas), 1), MAX_VIÑETAS_POR_SLIDE)]

    for idx, chunk in enumerate(chunks):
        slide = _blank_slide(prs)
        _set_slide_background(slide, COLOR_BG_DARK)
        _add_accent_bar(slide, prs)

        titulo_slide = titulo if idx == 0 else f"{titulo} (cont.)"
        _add_title_textbox(slide, titulo_slide, prs)
        _add_body_textbox(slide, chunk, prs)


def _build_referencias(prs: Presentation, refs: list):
    """
    Crea uno o más slides de Referencias al final, con MAX_REFS_POR_SLIDE por slide.
    """
    chunks = [refs[i:i+MAX_REFS_POR_SLIDE]
              for i in range(0, max(len(refs), 1), MAX_REFS_POR_SLIDE)]

    total = len(chunks)
    for idx, chunk in enumerate(chunks):
        slide = _blank_slide(prs)
        _set_slide_background(slide, COLOR_BG_REF)
        _add_accent_bar(slide, prs)

        # Título del slide de referencias
        label = "Referencias" if total == 1 else f"Referencias ({idx+1}/{total})"
        _add_title_textbox(slide, label, prs)

        # Cuerpo con las refs numeradas
        w, h = prs.slide_width, prs.slide_height
        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.6),
            w - Inches(1.1), h - Inches(2.1)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        start_n = idx * MAX_REFS_POR_SLIDE + 1
        for i, ref in enumerate(chunk):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = f"[{start_n + i}]  {ref}"
            run.font.name = FONT_BODY
            run.font.size = Pt(13)
            run.font.color.rgb = COLOR_REF_TEXT

        # Disclaimer en slide final de refs
        if idx == total - 1:
            disc = slide.shapes.add_textbox(
                Inches(0.6), h - Inches(0.55),
                prs.slide_width - Inches(1.2), Inches(0.4)
            )
            p_d = disc.text_frame.paragraphs[0]
            r_d = p_d.add_run()
            r_d.text = "⚠ Referencias generadas por IA. Verifique antes de publicar."
            r_d.font.size = Pt(10)
            r_d.font.italic = True
            r_d.font.color.rgb = RGBColor(0x88, 0x88, 0x99)


def _build_tabla_slide(prs: Presentation, titulo: str, headers: list, filas: list):
    """
    Crea uno o más slides con una tabla de datos cuando hay más de MAX_FILAS_TABLA_SLIDE filas.
    Encabezados en azul acento, filas alternadas claro/oscuro.
    """
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor

    COLOR_HEADER_BG  = RGBColor(0x4F, 0xA3, 0xFF)  # Azul acento para encabezados
    COLOR_HEADER_FG  = RGBColor(0xFF, 0xFF, 0xFF)  # Texto blanco en encabezado
    COLOR_ROW_ODD    = RGBColor(0x1E, 0x2A, 0x45)  # Fila impar - azul oscuro
    COLOR_ROW_EVEN   = RGBColor(0x15, 0x1E, 0x35)  # Fila par - aún más oscuro
    COLOR_ROW_TEXT   = RGBColor(0xD8, 0xD8, 0xE8)  # Texto de filas

    # Dividir filas en chunks si son muchas
    chunks = [filas[i:i+MAX_FILAS_TABLA_SLIDE]
              for i in range(0, max(len(filas), 1), MAX_FILAS_TABLA_SLIDE)]

    total = len(chunks)
    for idx, chunk in enumerate(chunks):
        slide = _blank_slide(prs)
        _set_slide_background(slide, COLOR_BG_DARK)
        _add_accent_bar(slide, prs)

        label = titulo if total == 1 else f"{titulo} ({idx+1}/{total})"
        _add_title_textbox(slide, label, prs)

        # Calcular dimensiones de la tabla
        w, h = prs.slide_width, prs.slide_height
        num_cols = len(headers) if headers else 1
        num_rows = len(chunk) + 1  # +1 para encabezados

        tbl_left   = Inches(0.4)
        tbl_top    = Inches(1.5)
        tbl_width  = w - Inches(0.8)
        tbl_height = h - Inches(2.0)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            tbl_left, tbl_top, tbl_width, tbl_height
        )
        table = table_shape.table

        # Ajustar columnas de ancho igual
        col_width = tbl_width // num_cols
        for col in table.columns:
            col.width = col_width

        def _set_cell(cell, texto, bold=False, bg_color=None, fg_color=COLOR_ROW_TEXT, font_size=13):
            cell.text = str(texto)
            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    run.font.name = FONT_BODY
                    run.font.color.rgb = fg_color
            if bg_color:
                from pptx.oxml.ns import qn
                from lxml import etree
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                # Eliminar relleno previo si existe
                for existing in tcPr.findall(qn('a:solidFill')):
                    tcPr.remove(existing)
                solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
                srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgbClr.set('val', str(bg_color))

        # Fila de encabezados
        for j, header in enumerate(headers[:num_cols]):
            _set_cell(table.cell(0, j), header, bold=True,
                      bg_color=COLOR_HEADER_BG, fg_color=COLOR_HEADER_FG, font_size=13)

        # Filas de datos
        for i, fila in enumerate(chunk):
            row_color = COLOR_ROW_ODD if i % 2 == 0 else COLOR_ROW_EVEN
            for j, valor in enumerate(fila[:num_cols]):
                _set_cell(table.cell(i + 1, j), valor,
                          bg_color=row_color, fg_color=COLOR_ROW_TEXT, font_size=12)


# ─────────────────────────────────────────────
# Funciones públicas
# ─────────────────────────────────────────────

def crear_ppt_base(filename: str, titulo: str, subtitulo: str, context: str = "outputs") -> str:
    """Crea una presentación PPTX básica con slide de portada."""
    try:
        prs = Presentation()
        _build_portada(prs, titulo, subtitulo)
        filepath = resolve_path(filename, context)
        prs.save(filepath)
        return f"PPT creado con éxito: {filepath}"
    except Exception as e:
        return f"Error creando PPT: {e}"


def actualizar_ppt(filename: str, reemplazos_dict: dict, context: str = "outputs") -> str:
    """Busca texto en todos los shapes y lo reemplaza con el diccionario."""
    try:
        filepath = resolve_path(filename, context)
        if not os.path.exists(filepath):
            return f"Error: No se encontró {filepath}"

        prs = Presentation(filepath)
        reemplazos = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in reemplazos_dict.items():
                            if key in run.text:
                                run.text = run.text.replace(key, value)
                                reemplazos += 1

        nuevo = os.path.splitext(filename)[0] + "_actualizado.pptx"
        prs.save(resolve_path(nuevo, context))
        return f"PPT actualizado con {reemplazos} reemplazos. Guardado: {nuevo}"
    except Exception as e:
        return f"Error actualizando PPT: {e}"


def insertar_texto_en_slide(filename: str, titulo: str, viñetas: list, context: str = "outputs") -> str:
    """Añade un nuevo slide con título y viñetas al final de la presentación."""
    try:
        filepath = resolve_path(filename, context)
        if not os.path.exists(filepath):
            return f"Error: No se encontró {filepath}"

        prs = Presentation(filepath)
        _build_contenido(prs, titulo, viñetas)
        prs.save(filepath)
        return f"Slide insertado con éxito en: {filepath}"
    except Exception as e:
        return f"Error insertando slide en PPT: {e}"


def crear_ppt_compleja_desde_json(filename: str, json_str: str, context: str = "outputs") -> str:
    """
    Crea una presentación completa desde un JSON estructurado.

    Tipos de slide soportados:
      portada     → titulo, subtitulo
      contenido   → titulo, viñetas (máx 5 por slide; divide automáticamente)
      cita        → texto de cita destacada
      indice      → titulo, viñetas (igual que contenido)
      referencias → items con las referencias APA (1 o más slides al final)
    """
    try:
        data = json.loads(json_str)
        if not data:
            return "Error: JSON vacío."

        filepath = resolve_path(filename, context)
        prs = Presentation()

        # Separar referencias para ponerlas al final
        slides_normales = []
        referencias_acumuladas = []

        for slide_data in data:
            tipo = slide_data.get("tipo", "contenido")
            if tipo == "referencias":
                referencias_acumuladas.extend(slide_data.get("items", []))
            else:
                slides_normales.append(slide_data)

        # Construir slides normales
        for slide_data in slides_normales:
            tipo  = slide_data.get("tipo", "contenido")
            titulo = slide_data.get("titulo", "")

            if tipo == "portada":
                _build_portada(prs, titulo, slide_data.get("subtitulo", ""))
            elif tipo == "cita":
                _build_cita(prs, slide_data.get("texto", titulo))
            elif tipo == "tabla_slide":
                headers = slide_data.get("headers", [])
                filas   = slide_data.get("filas", [])
                _build_tabla_slide(prs, titulo, headers, filas)
            else:  # contenido, indice
                viñetas = slide_data.get("viñetas", slide_data.get("items", []))
                _build_contenido(prs, titulo, viñetas)

        # Slides de referencias al final (si hay)
        if referencias_acumuladas:
            _build_referencias(prs, referencias_acumuladas)

        prs.save(filepath)
        ref_info = f" + {len(referencias_acumuladas)} referencias" if referencias_acumuladas else ""
        return f"PPT creado con {len(prs.slides)} slides{ref_info}: {filepath}"

    except json.JSONDecodeError as e:
        return f"Error parseando JSON de PPT: {e}\nJSON:\n{json_str[:500]}"
    except Exception as e:
        return f"Error creando PPT: {e}"
